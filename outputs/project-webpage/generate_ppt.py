"""
Idea2Top — 科技风 PPT 生成脚本
Usage: python generate_ppt.py
Output: Idea2Top_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# Color Palette (Tech Style)
# ============================================================
DARK_NAVY = RGBColor(0x0F, 0x19, 0x23)
NAVY = RGBColor(0x1A, 0x36, 0x5D)
CYAN = RGBColor(0x00, 0xA8, 0xFF)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1A, 0x20, 0x2C)
TEXT_MED = RGBColor(0x4A, 0x55, 0x68)
TEXT_LIGHT = RGBColor(0xA0, 0xAE, 0xC0)
GREEN = RGBColor(0x10, 0xB9, 0x81)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)

W = prs.slide_width
H = prs.slide_height

# ============================================================
# Helper Functions
# ============================================================

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text,
                 font_size=18, bold=False, color=TEXT_DARK,
                 alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_circle(slide, left, top, size, fill_color=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color=RGBColor(0x2A, 0x4A, 0x6A), width=Pt(2)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    connector.line.fill.solid()
    return connector

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

# Top accent bar
add_rect(slide, 0, 0, W, Pt(4), fill_color=CYAN)

# Decorative elements
add_rect(slide, 0, Emu(500000), Emu(80000), Emu(5500000), fill_color=RGBColor(0x15, 0x25, 0x35))
add_circle(slide, Emu(10500000), Emu(800000), Emu(220000), fill_color=RGBColor(0x00, 0x88, 0xCC))
add_circle(slide, Emu(11000000), Emu(1400000), Emu(120000), fill_color=PURPLE)
add_circle(slide, Emu(9800000), Emu(1200000), Emu(80000), fill_color=RGBColor(0x00, 0x88, 0x44))

# Title
add_text_box(slide, Emu(1000000), Emu(1800000), Emu(10000000), Emu(800000),
             "Idea2Top", font_size=60, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Emu(1000000), Emu(2600000), Emu(10000000), Emu(600000),
             "研究想法 → 顶刊实现", font_size=28, bold=False, color=CYAN, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Emu(1500000), Emu(3400000), Emu(9000000), Emu(400000),
             "将研究方法想法自动转化为顶刊级可复现实现的智能框架",
             font_size=16, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Emu(1500000), Emu(3800000), Emu(9000000), Emu(400000),
             "An intelligent framework for automated research methodology implementation",
             font_size=14, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Divider
add_rect(slide, Emu(4500000), Emu(4800000), Emu(3000000), Pt(1), fill_color=RGBColor(0x2A, 0x3A, 0x4A))

# Stats
stats_data = [("30+", "专业智能体"), ("50+", "算法知识库"), ("6", "核心流水线"), ("8", "领域覆盖")]
for i, (num, label) in enumerate(stats_data):
    x = Emu(1500000 + i * 2400000)
    add_text_box(slide, x, Emu(5100000), Emu(2200000), Emu(400000),
                 num, font_size=32, bold=True, color=CYAN, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Emu(5500000), Emu(2200000), Emu(300000),
                 label, font_size=12, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: About
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Side accent
add_rect(slide, 0, 0, Pt(6), H, fill_color=CYAN)

# Section label
add_rect(slide, 0, 0, Emu(2400000), Emu(400000), fill_color=NAVY)
add_text_box(slide, Emu(200000), Emu(60000), Emu(2000000), Emu(300000),
             "项目简介  ABOUT", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Left text
add_text_box(slide, Emu(800000), Emu(700000), Emu(5000000), Emu(500000),
             "Idea2Top 是什么？", font_size=28, bold=True, color=TEXT_DARK)

left_text = (
    "Idea2Top 是一个将研究方法想法自动转化为顶刊级\n"
    "可复现实现的智能框架。\n\n"
    "用户用自然语言描述研究想法，系统自动完成：\n"
    "• 文献对标 — 检索顶刊实现标准\n"
    "• 方案设计 — 数学建模与识别策略\n"
    "• 代码生成 — 顶刊级可运行实现\n"
    "• 实验验证 — 物理运行与对抗式QA\n"
    "• 交付 — 代码+图表+复现说明"
)
add_text_box(slide, Emu(800000), Emu(1300000), Emu(5000000), Emu(3200000),
             left_text, font_size=14, bold=False, color=TEXT_MED, alignment=PP_ALIGN.LEFT)

# Right side: 4 stat cards
stats = [
    ("30+", "专业智能体", "Specialized Agents"),
    ("50+", "算法知识库", "Algorithms Repository"),
    ("6", "核心流水线", "Core Workflows"),
    ("8", "领域覆盖", "Domain Coverage"),
]

for i, (num, label_zh, label_en) in enumerate(stats):
    row = i // 2
    col = i % 2
    x = Emu(6400000 + col * 2800000)
    y = Emu(800000 + row * 2400000)

    card = add_rounded_rect(slide, x, y, Emu(2500000), Emu(2000000),
                            fill_color=LIGHT_BG, border_color=CARD_BORDER)
    # Top accent
    accent_color = CYAN if i < 2 else PURPLE
    add_rect(slide, x, y, Emu(2500000), Pt(3), fill_color=accent_color)

    add_text_box(slide, x + Emu(100000), y + Emu(300000), Emu(2300000), Emu(600000),
                 num, font_size=40, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Emu(100000), y + Emu(900000), Emu(2300000), Emu(400000),
                 label_zh, font_size=16, bold=True, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Emu(100000), y + Emu(1300000), Emu(2300000), Emu(300000),
                 label_en, font_size=11, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)

# Header
add_rect(slide, 0, 0, W, Emu(500000), fill_color=NAVY)
add_text_box(slide, Emu(400000), Emu(100000), Emu(11000000), Emu(350000),
             "系统架构  ARCHITECTURE  —  三层多智能体分层协作",
             font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Level 1: Secretary
cy = Emu(800000)
card_w = Emu(2800000)
card_h = Emu(900000)
sec_x = Emu(4620000)

add_rounded_rect(slide, sec_x, cy, card_w, card_h, fill_color=NAVY, border_color=CYAN)
add_text_box(slide, sec_x + Emu(100000), cy + Emu(100000), card_w - Emu(200000), Emu(300000),
             "Secretary Agent", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, sec_x + Emu(100000), cy + Emu(380000), card_w - Emu(200000), Emu(400000),
             "秘书 Agent  |  任务分解守门人\n分解 → 确认 → 调度",
             font_size=10, bold=False, color=RGBColor(0xA0, 0xBE, 0xDD), alignment=PP_ALIGN.CENTER)

# Arrow
add_line(slide, sec_x + card_w // 2, cy + card_h, sec_x + card_w // 2, cy + card_h + Emu(300000))

# Level 2: Orchestrator
cy2 = cy + card_h + Emu(400000)
orch_x = Emu(4620000)
add_rounded_rect(slide, orch_x, cy2, card_w, card_h, fill_color=NAVY, border_color=PURPLE)
add_text_box(slide, orch_x + Emu(100000), cy2 + Emu(100000), card_w - Emu(200000), Emu(300000),
             "Orchestrator", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, orch_x + Emu(100000), cy2 + Emu(380000), card_w - Emu(200000), Emu(400000),
             "总协调人  |  意图识别 & 管线编排\n路由调度 → 上下文管理 → 质量门禁",
             font_size=10, bold=False, color=RGBColor(0xA0, 0xBE, 0xDD), alignment=PP_ALIGN.CENTER)

# Arrow
add_line(slide, orch_x + card_w // 2, cy2 + card_h, orch_x + card_w // 2, cy2 + card_h + Emu(300000))

# Horizontal distribution line
line_y = cy2 + card_h + Emu(350000)
add_line(slide, Emu(400000), line_y, Emu(11600000), line_y, color=RGBColor(0x3A, 0x5A, 0x7A), width=Pt(1))

# Level 3: Domain Agents
domains = [
    ("Literature", "文献检索与综述", 3),
    ("Topic Analysis", "前沿探测与选题", 3),
    ("Data & Viz", "数据清洗与可视化", 4),
    ("Experiment", "实验设计与优化", 3),
    ("Algorithm", "算法创造流水线", 5),
    ("Paper Format", "论文格式与排版", 3),
    ("Research QA", "科研问答与推导", 3),
    ("Kaggle", "竞赛端到端流水线", 7),
    ("Knowledge", "知识库管理", 1),
]

card_w2 = Emu(1180000)
card_h2 = Emu(1100000)

for i, (name, desc, count) in enumerate(domains):
    x = Emu(250000 + i * (card_w2.emu + 60000))
    y = line_y + Emu(100000)

    add_rounded_rect(slide, x, y, card_w2, card_h2, fill_color=WHITE, border_color=CARD_BORDER)
    add_rect(slide, x, y, card_w2, Pt(3), fill_color=CYAN)

    add_text_box(slide, x + Emu(50000), y + Emu(120000), card_w2 - Emu(100000), Emu(250000),
                 name, font_size=9, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Emu(50000), y + Emu(420000), card_w2 - Emu(100000), Emu(300000),
                 desc, font_size=7, bold=False, color=TEXT_MED, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Emu(50000), y + Emu(800000), card_w2 - Emu(100000), Emu(200000),
                 str(count) + " Agents", font_size=10, bold=True, color=CYAN, alignment=PP_ALIGN.CENTER)

    # Connection line
    add_line(slide, x + card_w2 // 2, line_y, x + card_w2 // 2, y,
             color=RGBColor(0x3A, 0x5A, 0x7A), width=Pt(1))

# ============================================================
# SLIDE 4: Core Workflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Header
add_rect(slide, 0, 0, W, Emu(500000), fill_color=NAVY)
add_text_box(slide, Emu(400000), Emu(100000), Emu(11000000), Emu(350000),
             "核心工作流  WORKFLOW  —  从想法到交付的完整流水线",
             font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

steps = [
    ("01", "解析与对标", "确认核心方法，检索顶刊\n标准实现规范", CYAN),
    ("02", "方案设计", "模型设定、识别策略\n估计方法、假设条件", RGBColor(0x00, 0x99, 0xCC)),
    ("03", "实现", "顶刊级可运行代码\n完整测试与使用示例", RGBColor(0x00, 0x88, 0xBB)),
    ("04", "实验与验证", "物理运行验证\n不信任AI口头报告", PURPLE),
    ("05", "对抗式QA", "Critic严格审查\n最多5轮循环", RGBColor(0x8B, 0x5C, 0xF6)),
    ("06", "交付", "代码+图表+方法\n复现说明 质量>=90", GREEN),
]

step_w = Emu(1750000)
step_h = Emu(3800000)
start_x = Emu(350000)
y_step = Emu(900000)

for i, (num, title, desc, color) in enumerate(steps):
    x = Emu(start_x.emu + i * (step_w.emu + 100000))
    add_rounded_rect(slide, x, y_step, step_w, step_h, fill_color=LIGHT_BG, border_color=CARD_BORDER)

    # Number circle
    circle_size = Emu(500000)
    circle_x = x + (step_w - circle_size) // 2
    add_circle(slide, circle_x, y_step + Emu(300000), circle_size, fill_color=color)
    add_text_box(slide, circle_x, y_step + Emu(380000), circle_size, Emu(400000),
                 num, font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, x + Emu(50000), y_step + Emu(1000000), step_w - Emu(100000), Emu(400000),
                 title, font_size=18, bold=True, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, x + Emu(80000), y_step + Emu(1500000), step_w - Emu(160000), Emu(1800000),
                 desc, font_size=11, bold=False, color=TEXT_MED, alignment=PP_ALIGN.CENTER)

    # Arrow connector
    if i < len(steps) - 1:
        arrow_x = x + step_w + Emu(20000)
        add_text_box(slide, arrow_x, y_step + Emu(1500000), Emu(80000), Emu(400000),
                     "→", font_size=22, bold=True, color=CYAN, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: Agent Cluster
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Header
add_rect(slide, 0, 0, W, Emu(500000), fill_color=NAVY)
add_text_box(slide, Emu(400000), Emu(100000), Emu(11000000), Emu(350000),
             "智能体集群  AGENT CLUSTER  —  30+ 专业智能体",
             font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

agents = [
    ("Secretary", "任务分解守门人"),
    ("Orchestrator", "总协调 & 编排"),
    ("Literature", "文献检索与综述"),
    ("Topic Analysis", "选题与前沿分析"),
    ("Data & Viz", "数据清洗与可视化"),
    ("Experiment", "实验设计与优化"),
    ("Algorithm", "算法创造管线"),
    ("Paper Format", "论文格式与排版"),
    ("Research QA", "科研问答与推导"),
    ("Kaggle", "竞赛流水线"),
    ("Knowledge", "知识库管理"),
]

card_w3 = Emu(2000000)
card_h3 = Emu(2200000)
cols = 4
start_x3 = Emu(350000)
y_step3 = Emu(800000)

for i, (name, desc) in enumerate(agents):
    col = i % cols
    row = i // cols
    x = Emu(start_x3.emu + col * (card_w3.emu + 120000))
    y = Emu(y_step3.emu + row * (card_h3.emu + 100000))

    add_rounded_rect(slide, x, y, card_w3, card_h3, fill_color=LIGHT_BG, border_color=CARD_BORDER)
    accent_color = CYAN if row == 0 else NAVY
    add_rect(slide, x, y, card_w3, Pt(3), fill_color=accent_color)

    add_text_box(slide, x + Emu(100000), y + Emu(200000), card_w3 - Emu(200000), Emu(400000),
                 name, font_size=14, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + Emu(100000), y + Emu(800000), card_w3 - Emu(200000), Emu(800000),
                 desc, font_size=10, bold=False, color=TEXT_MED, alignment=PP_ALIGN.LEFT)

# ============================================================
# SLIDE 6: Knowledge Repository
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Header
add_rect(slide, 0, 0, W, Emu(500000), fill_color=NAVY)
add_text_box(slide, Emu(400000), Emu(100000), Emu(11000000), Emu(350000),
             "知识库  KNOWLEDGE  —  50+ 算法与方法的知识沉淀",
             font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

categories = [
    ("因果推断", "Causal Inference", ["DID", "Causal Forest", "DML", "IV", "Conformal Pred."], CYAN),
    ("机器学习", "Machine Learning", ["XGBoost", "Random Forest", "SVM", "Clustering", "DR"], NAVY),
    ("时序与预测", "Time Series", ["ARIMA", "Prophet", "LSTM", "tsfresh", "Flow Matching"], RGBColor(0x00, 0x99, 0xCC)),
    ("深度学习", "Deep Learning", ["Autoencoder", "FlashAttention", "Enformer", "Geneformer", "AlphaFold3"], PURPLE),
    ("优化运筹", "Optimization", ["Adaptive ADMM", "Physarum Net", "Bayesian DRO", "Bilevel Opt."], RGBColor(0x8B, 0x5C, 0xF6)),
    ("贝叶斯统计", "Bayesian & Stats", ["Bayesian DRO", "Conformal Q", "Knockoffs", "Meta-Learners"], GREEN),
]

cat_w = Emu(3700000)
cat_h = Emu(2700000)

for i, (name_zh, name_en, items, color) in enumerate(categories):
    col = i % 3
    row = i // 3
    x = Emu(350000 + col * (cat_w.emu + 120000))
    y = Emu(800000 + row * (cat_h.emu + 100000))

    add_rounded_rect(slide, x, y, cat_w, cat_h, fill_color=LIGHT_BG, border_color=CARD_BORDER)
    add_rect(slide, x, y, cat_w, Pt(3), fill_color=color)

    add_text_box(slide, x + Emu(150000), y + Emu(200000), cat_w - Emu(300000), Emu(300000),
                 name_zh, font_size=18, bold=True, color=TEXT_DARK, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + Emu(150000), y + Emu(550000), cat_w - Emu(300000), Emu(250000),
                 name_en, font_size=11, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.LEFT)

    item_text = "\n".join(["•  " + item for item in items])
    add_text_box(slide, x + Emu(150000), y + Emu(1000000), cat_w - Emu(300000), Emu(1400000),
                 item_text, font_size=11, bold=False, color=TEXT_MED, alignment=PP_ALIGN.LEFT)

# ============================================================
# SLIDE 7: Quick Start Guide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

# Top bar
add_rect(slide, 0, 0, W, Pt(4), fill_color=CYAN)

add_text_box(slide, Emu(1000000), Emu(500000), Emu(10000000), Emu(500000),
             "使用指南  QUICK START", font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Emu(1000000), Emu(950000), Emu(10000000), Emu(300000),
             "三步上手  |  Get started in 3 steps",
             font_size=14, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

guide = [
    ("01", "提出研究想法", "Propose Your Idea",
     "用自然语言描述你的研究方法想法。\n例如：设计一个异质性处理效应的稳健估计量\n\n系统支持：研究想法 / 数据分析\n/ 竞赛任务 / 文献综述"),
    ("02", "确认分解方案", "Review the Plan",
     "秘书Agent自动分析任务并输出分解方案。\n检查子任务列表，确认工具/配色/规模/格式。\n\n确认后 Orchestrator 开始调度执行。"),
    ("03", "获取顶刊级交付", "Get Your Delivery",
     "系统自动完成：文献对标 → 方案设计\n→ 代码生成 → 实验验证 → 对抗式QA\n\n交付：可运行代码+图表+复现说明"),
]

for i, (num, title_zh, title_en, desc) in enumerate(guide):
    x = Emu(600000 + i * 3700000)
    y = Emu(1600000)
    sw = Emu(3400000)
    sh = Emu(4000000)

    add_rounded_rect(slide, x, y, sw, sh,
                     fill_color=RGBColor(0x15, 0x25, 0x38),
                     border_color=RGBColor(0x2A, 0x3A, 0x5A))

    # Step number circle
    csize = Emu(600000)
    add_circle(slide, x + (sw - csize) // 2, y + Emu(300000), csize, fill_color=CYAN)
    add_text_box(slide, x + (sw - csize) // 2, y + Emu(380000), csize, Emu(400000),
                 num, font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Titles
    add_text_box(slide, x + Emu(150000), y + Emu(1200000), sw - Emu(300000), Emu(400000),
                 title_zh, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Emu(150000), y + Emu(1600000), sw - Emu(300000), Emu(300000),
                 title_en, font_size=10, bold=False, color=CYAN, alignment=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, x + Emu(500000), y + Emu(2000000), sw - Emu(1000000), Pt(1),
             fill_color=RGBColor(0x2A, 0x3A, 0x5A))

    # Description
    add_text_box(slide, x + Emu(200000), y + Emu(2200000), sw - Emu(400000), Emu(1500000),
                 desc, font_size=11, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide, Emu(1500000), Emu(5800000), Emu(9000000), Emu(500000),
             "只需用自然语言描述，剩下的交给 Agent 集群",
             font_size=12, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

add_rect(slide, 0, 0, W, Pt(4), fill_color=CYAN)
add_rect(slide, Emu(500000), Emu(2500000), Emu(11000000), Pt(1), fill_color=RGBColor(0x2A, 0x3A, 0x5A))

add_text_box(slide, Emu(1000000), Emu(2200000), Emu(10000000), Emu(800000),
             "Thank You", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Emu(1000000), Emu(3000000), Emu(10000000), Emu(600000),
             "Idea2Top  —  研究想法 → 顶刊实现",
             font_size=20, bold=False, color=CYAN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Emu(2500000), Emu(3800000), Emu(7000000), Emu(400000),
             "数学建模半自动 · Research Assistant Framework",
             font_size=14, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Emu(2500000), Emu(4200000), Emu(7000000), Emu(400000),
             "Powered by Claude Multi-Agent Architecture",
             font_size=12, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_rect(slide, Emu(500000), Emu(5200000), Emu(11000000), Pt(1), fill_color=RGBColor(0x2A, 0x3A, 0x5A))
add_text_box(slide, Emu(1000000), Emu(5400000), Emu(10000000), Emu(400000),
             "© 2026 Idea2Top Research Framework",
             font_size=11, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = "E:/wuyi/数学建模半自动/research-assistant/outputs/project-webpage/Idea2Top_Presentation.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("Done!")
